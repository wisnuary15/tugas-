from pptx import Presentation
from pptx.util import Inches

# Buat presentasi
prs = Presentation()

# Slide 1: Judul
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
title.text = "Evaluasi Kinerja Keuangan Perusahaan Teknologi di Bursa Efek Istanbul Menggunakan Metode TOPSIS"

# Slide 2: Pendahuluan
slide_intro = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_intro.shapes.title, slide_intro.shapes.placeholders[1]
title.text = "Pendahuluan"
content.text = ("Teknik TOPSIS digunakan untuk mengevaluasi kinerja keuangan perusahaan teknologi.\n"
                "Perusahaan yang dievaluasi adalah perusahaan yang terdaftar di Bursa Efek Istanbul dari tahun 2009-2011.\n"
                "10 rasio keuangan digunakan sebagai kriteria evaluasi.")

# Slide 3: Metodologi TOPSIS
slide_method = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_method.shapes.title, slide_method.shapes.placeholders[1]
title.text = "Metodologi TOPSIS"
content.text = ("1. Normalisasi Data\n"
                "2. Menentukan Solusi Ideal dan Anti-Ideal\n"
                "3. Menghitung Jarak Pemisahan\n"
                "4. Menghitung Kedekatan Relatif ke Solusi Ideal (C*)\n"
                "5. Membuat Peringkat")

# Slide 4: Data Asli (2009)
slide_data = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_data.shapes.title, slide_data.shapes.placeholders[1]
title.text = "Data Asli (2009)"
content.text = "Tabel data asli dapat dilihat pada file CSV yang disertakan."

# Slide 5: Matriks Normalisasi (2009)
slide_norm = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_norm.shapes.title, slide_norm.shapes.placeholders[1]
title.text = "Matriks Normalisasi (2009)"
content.text = "Tabel matriks normalisasi dapat dilihat pada file CSV yang disertakan."

# Slide 6: Peringkat Perusahaan (2009)
slide_rank = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_rank.shapes.title, slide_rank.shapes.placeholders[1]
title.text = "Peringkat Perusahaan (2009)"
content.text = "Peringkat perusahaan berdasarkan metode TOPSIS dapat dilihat pada tabel berikut:\n"
rank_table = norm_df[['Perusahaan', 'C*', 'Peringkat']]
content.text += rank_table.to_string(index=False)

# Simpan presentasi
pptx_file = '/mnt/data/Evaluasi_Kinerja_TOPSIS.pptx'
prs.save(pptx_file)

pptx_file &#8203;:citation[oaicite:0]{index=0}&#8203;
